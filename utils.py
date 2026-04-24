import os
from docx import Document
import openpyxl
from pptx import Presentation
from PIL import Image
import pytesseract
import PyPDF2
from pdf2image import convert_from_path

pytesseract.pytesseract.tesseract_cmd = "/usr/bin/tesseract"
POPPLER_PATH = "/usr/bin"

def extract_text_from_pdf(pdf_path):
    text = ""

    try:
        with open(pdf_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    except Exception as e:
        print("PDF extract error:", e)

    if len(text.strip()) < 50:
        pages = convert_from_path(pdf_path, poppler_path=POPPLER_PATH)
        for page in pages:
            text += pytesseract.image_to_string(page, lang="fra+eng")

    return text

def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)

    elif ext == ".docx":
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + " "

    elif ext in [".xlsx", ".xls"]:
        wb = openpyxl.load_workbook(file_path)
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell:
                        text += str(cell) + " "

    elif ext == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "

    elif ext in [".png", ".jpg", ".jpeg"]:
        img = Image.open(file_path)
        text += pytesseract.image_to_string(img, lang="fra+eng")

    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

    return text

def smart_extract(text):
    words = text.split()
    if len(words) < 600:
        return text
    return " ".join(words[:300] + words[-300:])

def summarize_text(text, summarizer, chunk_size=800):
    summaries = []

    for i in range(0, len(text), chunk_size):
        chunk = text[i:i+chunk_size]

        summary = summarizer(
            "summarize: " + chunk,
            max_length=100,
            min_length=30,
            do_sample=False
        )

        summaries.append(summary[0]["summary_text"])

    return " ".join(summaries)